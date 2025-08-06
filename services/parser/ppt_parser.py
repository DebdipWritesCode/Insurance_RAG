import os
import tempfile
from pptx import Presentation
from pdf2image import convert_from_path
import pytesseract

import subprocess

def pptx_to_pdf(pptx_path: str, output_pdf_path: str) -> None:
    output_dir = os.path.dirname(output_pdf_path)
    command = [
        "soffice",
        "--headless",
        "--convert-to", "pdf",
        "--outdir", output_dir,
        pptx_path
    ]
    try:
        result = subprocess.run(command, capture_output=True, check=True)
        print("✅ LibreOffice Output:", result.stdout.decode())
    except subprocess.CalledProcessError as e:
        print("❌ LibreOffice failed:", e.stderr.decode())
        raise RuntimeError("PDF conversion failed") from e
    
def extract_text_from_pdf_ocr(pdf_path: str) -> list[str]:
    slides = convert_from_path(pdf_path, dpi=300)
    all_chunks = []

    for i, image in enumerate(slides):
        print(f"🔍 OCR on slide {i+1}")
        text = pytesseract.image_to_string(image)
        chunks = [chunk.strip() for chunk in text.split("\n\n") if chunk.strip()]
        all_chunks.extend(chunks)

    print(f"🖼️ Extracted {len(all_chunks)} chunks from OCR (PDF-based)")
    return all_chunks

def extract_text_from_pptx(pptx_path: str) -> list[str]:
    prs = Presentation(pptx_path)
    full_text = []
    slides_needing_ocr = []

    for i, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        
        if slide_text:
            full_text.append("\n".join(slide_text))
        else:
            print(f"⚠️ No text found on slide {i+1}, marking for OCR.")
            slides_needing_ocr.append(i)

    if not slides_needing_ocr:
        print("📊 All text extracted from PPTX without OCR.")
        return [chunk.strip() for chunk in full_text if chunk.strip()]

    with tempfile.TemporaryDirectory() as tmpdir:
        pdf_path = os.path.join(tmpdir, "converted.pdf")
        pptx_to_pdf(pptx_path, pdf_path)
        all_images = convert_from_path(pdf_path, dpi=300)

        for i in slides_needing_ocr:
            image = all_images[i]
            print(f"🔍 OCR on slide {i+1}")
            ocr_text = pytesseract.image_to_string(image)
            ocr_chunks = [chunk.strip() for chunk in ocr_text.split("\n\n") if chunk.strip()]
            full_text.extend(ocr_chunks)

    print(f"📊 Extracted {len(full_text)} chunks from PPTX (with OCR fallback)")
    print(f"🔍 First chunk (100 chars): {repr(full_text[0][:100]) if full_text else 'No chunks'}")

    return [chunk.strip() for chunk in full_text if chunk.strip()]
